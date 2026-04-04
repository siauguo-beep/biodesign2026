from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BG = RGBColor(0x14, 0x12, 0x10)
BG2 = RGBColor(0x1E, 0x1A, 0x17)
FG = RGBColor(0xE6, 0xDF, 0xD6)
MUTED = RGBColor(0x9A, 0x8F, 0x84)
ACCENT = RGBColor(0xC4, 0x5C, 0x2A)
SIGNAL = RGBColor(0x5C, 0xB8, 0x9A)
CHROME = RGBColor(0x6B, 0x7A, 0x8F)


TITLE = "Extinction Archive"
SUBTITLE = "Umwelt hypothesis dossiers"
FOOTER = "Extinction Archive · Biodigital chronobiology · BDC 2026"


SLIDES = [
    {
        "type": "lead",
        "title": "Extinction Archive",
        "subtitle": "Umwelt hypothesis dossiers",
        "tagline": "circadian · migration · synchrony",
        "body": [
            "Biodesign Challenge 2026 · Convergent Life",
            "Digital Art & AI Technology — MDes",
            "Literature-grounded memorial · biodigital experience",
        ],
    },
    {
        "type": "bullets",
        "eyebrow": "Why this project now",
        "title": "Extinction removes more than bodies.",
        "bullets": [
            "It erases migration pulses, social synchrony, and seasonal fit to landscape.",
            "Public memory often normalizes ecological absence after enough time passes.",
            "This project asks how constrained computation can make loss legible without pretending we fully know another species' Umwelt.",
        ],
        "callout_title": "Framing question",
        "callout_body": "How can an archive make extinction felt as both biological disappearance and cultural amnesia?",
    },
    {
        "type": "bullets",
        "eyebrow": "Design proposition",
        "title": "Umwelt as memory technology",
        "bullets": [
            "Primary frame: bounded sensory-temporal hypothesis space.",
            "Secondary frame: collective memory fracture and ethical remembrance.",
            "Visitors move through Memory Sites, Species Dossier, Polyphony Mixer, and Ethics Console.",
            "The archive is not resurrection; uncertainty remains visible at every stage.",
        ],
        "callout_title": "Tone",
        "callout_body": "Museum sobriety overall. Awe appears briefly, mainly inside the mixer.",
    },
    {
        "type": "columns",
        "eyebrow": "Core architecture",
        "title": "Three load-bearing modules",
        "left_title": "Species Dossier",
        "left_body": [
            "Evidence cards",
            "Temporal niche",
            "Sensory proxies",
            "Extinction mechanism",
            "360° constrained scene",
        ],
        "mid_title": "Polyphony Mixer",
        "mid_body": [
            "Seasonal layers",
            "Diel rhythm",
            "Density and synchrony",
            "Audible coexistence",
            "Collapse as retuning",
        ],
        "right_title": "Ethics Console",
        "right_body": [
            "Evidence sorting",
            "Forced tradeoffs",
            "Rule-based outcomes",
            "Sonic finale shifts",
            "Action over nostalgia",
        ],
    },
    {
        "type": "columns",
        "eyebrow": "Hero species",
        "title": "Mammoth-first, pigeon as structural echo",
        "left_title": "Woolly mammoth",
        "left_body": [
            "Tusk isotope trails",
            "Arctic photoperiod",
            "Comparative genomics",
            "Ground-level pacing",
            "Slow seasonal time",
        ],
        "mid_title": "Passenger pigeon",
        "mid_body": [
            "Mass flock synchrony",
            "Density thresholds",
            "Social collapse",
            "Sky-darkening scale",
            "Thinness as silence",
        ],
        "right_title": "Memory sites MVP",
        "right_body": [
            "2-5 curated pins",
            "Last ranges and dates",
            "Single-species deep links",
            "No global atlas in MVP",
        ],
    },
    {
        "type": "bullets",
        "eyebrow": "Biological layer",
        "title": "Only traceable evidence enters the archive",
        "bullets": [
            "Migration isotopes, morphology as sensory proxy, comparative genomics, historical range data, and acoustic ecology of relatives.",
            "Smell cues remain metaphorical or synesthetic only and are always labeled Speculative.",
            "Every strong line in the UI is tagged: Cited, Modeled, or Speculative.",
        ],
        "callout_title": "Rule",
        "callout_body": "If the biology can be removed and the piece still works, the project has failed its biodigital test.",
    },
    {
        "type": "bullets",
        "eyebrow": "Digital layer",
        "title": "Expressive, but never unconstrained",
        "bullets": [
            "Primary image strategy: diffusion plus structural conditioning.",
            "Support layer: procedural shaders for light, fog, atmosphere, and temporal credibility.",
            "A-Frame-first WebXR with desktop fallback and mobile WebAR entry points.",
            "Web Audio DSP ties season, density, and time-of-day to the sonic field.",
            "R and ggplot provide auditable scientific visuals for dossiers and appendix slides.",
        ],
        "callout_title": "Tooling",
        "callout_body": "A-Frame first, Three.js as needed inside the same XR ecosystem.",
    },
    {
        "type": "bullets",
        "eyebrow": "Physical contrast",
        "title": "Archive object + living bench",
        "bullets": [
            "3D-printed focal object anchors the body of the extinct species in gallery space.",
            "Mobile WebAR opens dossier content from the physical object.",
            "Safe tactile materials symbolize tundra and steppe without claiming literal reconstruction.",
            "Plants plus soil moisture, light, and temperature sensors drive one live audio bus.",
        ],
        "callout_title": "Installation label",
        "callout_body": "Not a simulation of extinct habitat — an index of present life.",
    },
    {
        "type": "bullets",
        "eyebrow": "Ethics and reflection",
        "title": "Reflection is part of the interface",
        "bullets": [
            "Digital resurrection is kept distinct from laboratory de-extinction discourse.",
            "Visitors face choices around conservation resources, land-use pressure, and moral hazard.",
            "AI-generated confidence is actively countered through tier labels, provenance, and rule-based consequences.",
            "The project avoids speaking for communities it has not collaborated with.",
        ],
        "callout_title": "Desired effect",
        "callout_body": "Move from grief to responsibility, not from grief to spectacle.",
    },
    {
        "type": "bullets",
        "eyebrow": "BDC alignment",
        "title": "Why this fits Biodesign Challenge 2026",
        "bullets": [
            "Narrative: a clear emotional and spatial journey supported by physical and digital prototypes.",
            "Concept: Umwelt plus temporal niche plus collective memory creates a distinct biodigital proposition.",
            "Context: biodiversity loss, de-extinction ethics, misinformation risk, and sustainable installation logic.",
            "Reflection: transparent iteration, limits, evidence tiers, and expert-reviewable claims.",
        ],
        "callout_title": "Prize target",
        "callout_body": "Primary target: BDC Prize for Biodigital Excellence.",
    },
    {
        "type": "timeline",
        "eyebrow": "1-month vertical slice",
        "title": "Build the hero demo first",
        "steps": [
            ("Week 1", "Lock evidence cards, dossier fields, memory sites, ethics branches, and consultant outreach."),
            ("Week 2", "Ship a mammoth vertical slice: dossier, mixer layer, A-Frame scene, and film-track kickoff."),
            ("Week 3", "Integrate passenger pigeon, ethics logic, WebAR, tactile table, and living sensor pipeline."),
            ("Week 4", "Polish, rehearse, finish trailer, finalize appendix citations, and harden demo fallback."),
        ],
        "callout_title": "Parallel rule",
        "callout_body": "The film track starts after Week 2 and must not block the core interaction milestone.",
    },
    {
        "type": "bullets",
        "eyebrow": "Deliverables map",
        "title": "What the final package includes",
        "bullets": [
            "English planning document ready for the repo.",
            "Editable PowerPoint deck in project theme.",
            "GitHub slide README with Office preview link and PDF preview link.",
            "A methods appendix culture: citations and limits developed in parallel, not at the end.",
        ],
        "callout_title": "Closing question",
        "callout_body": "What should responsibility sound like after a species has already disappeared?",
    },
]


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(9.9),
        Inches(0.18),
        Inches(3.43),
        prs.slide_height - Inches(0.18),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = BG2
    band.fill.transparency = 0.15
    band.line.fill.background()


def add_footer(slide, idx):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(11.9), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{FOOTER}  {idx} / {len(SLIDES)}"
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.LEFT


def add_eyebrow(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(4.8), Inches(0.32))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = ACCENT


def add_title(slide, text, top=0.85, width=7.7):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(width), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = FG


def add_bullet_box(slide, items, left=0.8, top=2.0, width=7.0, height=3.8, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Helvetica Neue"
        p.font.size = Pt(font_size)
        p.font.color.rgb = FG
        p.space_after = Pt(10)
        p.bullet = True


def add_callout(slide, title, body):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(1.7), Inches(3.5), Inches(3.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG2
    shape.line.color.rgb = CHROME

    title_box = slide.shapes.add_textbox(Inches(8.75), Inches(2.0), Inches(2.8), Inches(0.4))
    tp = title_box.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title.upper()
    tr.font.name = "Helvetica Neue"
    tr.font.size = Pt(12)
    tr.font.bold = True
    tr.font.color.rgb = SIGNAL

    body_box = slide.shapes.add_textbox(Inches(8.75), Inches(2.45), Inches(2.8), Inches(2.4))
    bp = body_box.text_frame.paragraphs[0]
    br = bp.add_run()
    br.text = body
    br.font.name = "Helvetica Neue"
    br.font.size = Pt(16)
    br.font.color.rgb = FG


def add_column_group(slide, left, title, items, box_width=3.45):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.95), Inches(box_width), Inches(3.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG2
    shape.line.color.rgb = CHROME

    title_box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(2.12), Inches(box_width - 0.36), Inches(0.45))
    tp = title_box.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.name = "Helvetica Neue"
    tr.font.size = Pt(18)
    tr.font.bold = True
    tr.font.color.rgb = FG

    add_bullet_box(slide, items, left=left + 0.18, top=2.65, width=box_width - 0.36, height=2.9, font_size=14)


def add_timeline(slide, steps):
    y = 2.0
    for i, (week, desc) in enumerate(steps):
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(y), Inches(7.1), Inches(0.85)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BG2
        line.line.color.rgb = CHROME

        wk = slide.shapes.add_textbox(Inches(1.1), Inches(y + 0.17), Inches(1.0), Inches(0.3))
        wp = wk.text_frame.paragraphs[0]
        wr = wp.add_run()
        wr.text = week
        wr.font.name = "Helvetica Neue"
        wr.font.size = Pt(15)
        wr.font.bold = True
        wr.font.color.rgb = ACCENT

        ds = slide.shapes.add_textbox(Inches(2.05), Inches(y + 0.12), Inches(5.6), Inches(0.48))
        dp = ds.text_frame.paragraphs[0]
        dr = dp.add_run()
        dr.text = desc
        dr.font.name = "Helvetica Neue"
        dr.font.size = Pt(14)
        dr.font.color.rgb = FG
        y += 1.0


def render_lead(slide, data):
    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(8.2), Inches(0.9))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = data["title"]
    r.font.name = "Helvetica Neue"
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = FG

    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(1.95), Inches(8.6), Inches(0.7))
    sp = subtitle.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = data["subtitle"]
    sr.font.name = "Helvetica Neue"
    sr.font.size = Pt(24)
    sr.font.color.rgb = FG

    tag = slide.shapes.add_textbox(Inches(0.82), Inches(3.0), Inches(5.0), Inches(0.4))
    tp = tag.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = data["tagline"]
    tr.font.name = "Helvetica Neue"
    tr.font.size = Pt(13)
    tr.font.bold = True
    tr.font.color.rgb = ACCENT

    body = slide.shapes.add_textbox(Inches(0.82), Inches(3.7), Inches(6.6), Inches(1.6))
    tf = body.text_frame
    for i, line in enumerate(data["body"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Helvetica Neue"
        p.font.size = Pt(18)
        p.font.color.rgb = FG
        p.space_after = Pt(8)

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.35), Inches(1.2), Inches(3.2), Inches(4.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = BG2
    box.line.color.rgb = CHROME

    q = slide.shapes.add_textbox(Inches(8.7), Inches(1.7), Inches(2.5), Inches(3.5))
    qp = q.text_frame.paragraphs[0]
    qr = qp.add_run()
    qr.text = "A biodigital memorial where extinction is experienced as rhythm, silence, and consequence."
    qr.font.name = "Helvetica Neue"
    qr.font.size = Pt(22)
    qr.font.color.rgb = FG


def render_bullets(slide, data):
    add_eyebrow(slide, data["eyebrow"])
    add_title(slide, data["title"])
    add_bullet_box(slide, data["bullets"])
    add_callout(slide, data["callout_title"], data["callout_body"])


def render_columns(slide, data):
    add_eyebrow(slide, data["eyebrow"])
    add_title(slide, data["title"], width=10.5)
    add_column_group(slide, 0.7, data["left_title"], data["left_body"])
    add_column_group(slide, 4.25, data["mid_title"], data["mid_body"])
    add_column_group(slide, 7.8, data["right_title"], data["right_body"], box_width=3.9)


def render_timeline(slide, data):
    add_eyebrow(slide, data["eyebrow"])
    add_title(slide, data["title"])
    add_timeline(slide, data["steps"])
    add_callout(slide, data["callout_title"], data["callout_body"])


RENDERERS = {
    "lead": render_lead,
    "bullets": render_bullets,
    "columns": render_columns,
    "timeline": render_timeline,
}


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for idx, spec in enumerate(SLIDES, start=1):
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    RENDERERS[spec["type"]](slide, spec)
    add_footer(slide, idx)


out = Path("slides/export/BDC_Extinction_Archive_EN.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"Saved {out}")
