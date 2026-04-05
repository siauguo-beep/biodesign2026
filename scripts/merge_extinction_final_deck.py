#!/usr/bin/env python3
"""
Build the final ≤25-slide deck (PPTX only):
- Part I: BDC_Extinction_Archive_EN.pptx (theme, backgrounds, in-place text fixes)
- Part II: Curated literature / methods slides from BDC_Summit_Extinction_Archive_2026.pptx
  using the same "Title and Content" layout as the template (consistent master background)
- Closing: Thank you

Run: .venv/bin/python scripts/merge_extinction_final_deck.py
"""

from __future__ import annotations

import re
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent.parent
EXPORT = ROOT / "slides" / "export"
TEMPLATE = EXPORT / "BDC_Extinction_Archive_EN.pptx"
SUMMIT = EXPORT / "BDC_Summit_Extinction_Archive_2026.pptx"

# PPT filename (display / export) — user-specified naming
OUT = (
    EXPORT
    / "{Extinction Archive} Umwelt Hypothesis Dossiers——AI memorial for lost species · Sensory time capsule.pptx"
)

TITLE_LINE1 = "{Extinction Archive} Umwelt Hypothesis Dossiers"
TITLE_LINE2 = "AI memorial for lost species · Sensory time capsule"

# Hard cap: 12 (template) + len(SUMMIT_APPEND_INDICES) + 1 (thank you) ≤ 25
MAX_TOTAL_SLIDES = 25

# Curated summit slides: concept, scenes, pipeline, colonial/ethics, literature, rubric, repo
SUMMIT_APPEND_INDICES = [
    6,   # Core design question
    7,   # Bio × digital quadrants
    8,   # Two dossiers
    9,   # Scene IDs
    10,  # Design techniques
    12,  # Four-week spine
    13,  # Colonial / Country
    14,  # Sherkow & Greely
    15,  # Literature — mammoth
    16,  # Literature — thylacine
    17,  # Rubric
    19,  # Repository / pointers
]


def replace_exact_paragraph(slide, exact: str, new: str) -> None:
    """Replace only when paragraph text equals `exact` (avoids substring collisions)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = "".join(run.text for run in para.runs)
            if full.strip() != exact.strip():
                continue
            for i, run in enumerate(para.runs):
                run.text = new if i == 0 else ""


def replace_in_all_text_shapes(slide, mapping: list[tuple[str, str]]) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = "".join(run.text for run in para.runs)
            if not full.strip():
                continue
            new_full = full
            for old, new in mapping:
                if old in new_full:
                    new_full = new_full.replace(old, new)
            if new_full != full:
                for i, run in enumerate(para.runs):
                    run.text = new_full if i == 0 else ""


def slide_text_blob(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            parts.append(shape.text.strip())
    return "\n\n".join(parts)


def add_title_content_slide(prs: Presentation, title: str, body_lines: list[str]) -> None:
    """Uses template layout 1 — same slide master / background as narrative slides."""
    blank_slide = prs.slides.add_slide(prs.slide_layouts[1])
    if blank_slide.shapes.title:
        blank_slide.shapes.title.text = title
    ph = blank_slide.placeholders[1]
    tf = ph.text_frame
    tf.clear()
    first = True
    for line in body_lines:
        line = line.strip()
        if not line:
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(14)


def add_thank_you_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Thank you"
    lines = [
        TITLE_LINE1,
        TITLE_LINE2,
        "Macau University · Biodesign Challenge 2026 · Biodigital Excellence",
        "GUO XIAO YUE (MC569254) · LIU JIA QUN (MC569293) · Mentor: Atticus SIMS",
    ]
    ph = slide.placeholders[1]
    tf = ph.text_frame
    tf.clear()
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(16)


def renumber_footers(prs: Presentation) -> None:
    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full = "".join(run.text for run in para.runs)
                if "BDC 2026" not in full or "/" not in full:
                    continue
                new_full = re.sub(r"\d+\s*/\s*\d+", f"{idx + 1} / {total}", full)
                if new_full != full:
                    for i, run in enumerate(para.runs):
                        run.text = new_full if i == 0 else ""


def merge() -> None:
    shutil.copy2(TEMPLATE, OUT)
    prs = Presentation(str(OUT))

    # Slide 1 — long title (avoid replacing short substring before full footer line)
    replace_in_all_text_shapes(
        prs.slides[0],
        [
            (
                "Extinction Archive · Biodigital chronobiology · BDC 2026  1 / 12",
                f"{TITLE_LINE1} · Biodigital chronobiology · BDC 2026  1 / 12",
            ),
            ("Umwelt hypothesis dossiers", TITLE_LINE2),
            (
                "circadian · migration · synchrony",
                "circadian · migration · synchrony · multisensory memory",
            ),
            (
                "Literature-grounded memorial · biodigital experience",
                "Literature-grounded memorial · biodigital experience · memory reshaping\n"
                "Macau University · Mentor: Atticus SIMS · GUO XIAO YUE (MC569254) · LIU JIA QUN (MC569293)",
            ),
            (
                "A biodigital memorial where extinction is experienced as rhythm, silence, and consequence.",
                "A biodigital memorial where extinction is experienced as rhythm, multisensory reminiscence, and consequence — with cognitive limits and uncertainty kept visible.",
            ),
        ],
    )
    replace_exact_paragraph(prs.slides[0], "Extinction Archive", TITLE_LINE1)

    # Slide 4 — Core architecture: memory reshaping + multisensory reminiscence
    replace_in_all_text_shapes(
        prs.slides[3],
        [
            (
                "Three load-bearing modules",
                "Memory reshaping (foreground): multisensory reminiscence retrains attention on extinct Umwelten; guardrails address collective cognitive impairment, flattening, and misremembering.\nThree load-bearing modules",
            ),
            ("Evidence cards", "Evidence cards · epistemic tiers"),
            ("Temporal niche", "Temporal niche · diel / seasonal anchors"),
            ("Sensory proxies", "Sensory proxies · multisensory reminiscence paths"),
            (
                "Extinction mechanism",
                "Extinction mechanism · colonial + ecological chronologies",
            ),
            (
                "Action over nostalgia",
                "Action over nostalgia\nMixer + Console: reminiscence without false certainty",
            ),
        ],
    )

    replace_in_all_text_shapes(
        prs.slides[4],
        [
            (
                "Mammoth-first, pigeon as structural echo",
                "Memory reshaping through paired hero dossiers\nMammoth + thylacine dossiers (two deep case studies)",
            ),
            (
                "Woolly mammoth",
                "Woolly mammoth — seasonal Umwelt for slow, somatic reminiscence",
            ),
            ("Passenger pigeon", "Thylacine (Tasmanian tiger)"),
            ("Mass flock synchrony", "Orbit / diel activity (Pozniak 2018)"),
            (
                "Density thresholds",
                "POV metaphor via RGC models (Mass & Supin 2020, Interpolated)",
            ),
            (
                "Social collapse",
                "Colonial extinction history (Paddle; Sleightholme & Campbell)",
            ),
            (
                "Sky-darkening scale",
                "Palawa / Country framing (Rimmer; Lehman; Clements; Schlunke)",
            ),
            (
                "Thinness as silence",
                "Ethics Console + reflection log — cognition-friendly branching",
            ),
        ],
    )

    replace_in_all_text_shapes(
        prs.slides[7],
        [
            ("Archive object + living bench", "Minimal physical anchor + WebXR entry (Summit build)"),
            (
                "3D-printed focal object anchors the body of the extinct species in gallery space.\nMobile WebAR opens dossier content from the physical object.\nSafe tactile materials symbolize tundra and steppe without claiming literal reconstruction.\nPlants plus soil moisture, light, and temperature sensors drive one live audio bus.",
                "3D-printed or surrogate focal object anchors attention in gallery space.\nQR / URL opens the WebXR memorial (A-Frame / Three.js).\nWeb-first scope: no living biosensor demo in v1; one citation one-sheet for judges.\nOptional WebAR from the printed anchor.",
            ),
            (
                "Not a simulation of extinct habitat — an index of present life.",
                "Not a simulation of extinct habitat — an index of evidence, uncertainty, and responsibility.",
            ),
        ],
    )

    replace_in_all_text_shapes(
        prs.slides[10],
        [
            (
                "Build the hero demo first",
                "Build the hero demo first — memory reshaping in the loop",
            ),
            (
                "Lock evidence cards, dossier fields, memory sites, ethics branches, and consultant outreach.",
                "Lock evidence cards, dossier fields, multisensory reminiscence beats, memory sites, ethics branches, plain-language / cognitive-access patterns, consultant outreach.",
            ),
            (
                "Ship a mammoth vertical slice: dossier, mixer layer, A-Frame scene, and film-track kickoff.",
                "Ship a mammoth vertical slice: dossier, polyphonic mixer reminiscence layers, WebXR scene, film-track kickoff.",
            ),
            (
                "Integrate passenger pigeon, ethics logic, WebAR, tactile table, and living sensor pipeline.",
                "Integrate thylacine dossier, Indigenous context UI, ethics fork (Sherkow & Greely), reflection log, WebXR polish — rehearse multisensory pacing + confusion testing.",
            ),
            (
                "Polish, rehearse, finish trailer, finalize appendix citations, and harden demo fallback.",
                "Polish, rehearse, finish trailer, finalize appendix citations, harden demo fallback, "
                "validate reminiscence overload safeguards.",
            ),
            (
                "The film track starts after Week 2 and must not block the core interaction milestone.",
                "The film track starts after Week 2 and must not block the core interaction milestone.\nParallel rule (memory): remix and polyphony remain legible under fatigue — no sentimental erasure of uncertainty.",
            ),
        ],
    )

    summit_prs = Presentation(str(SUMMIT))
    for idx in SUMMIT_APPEND_INDICES:
        if idx >= len(summit_prs.slides):
            continue
        blob = slide_text_blob(summit_prs.slides[idx])
        if not blob.strip():
            continue
        lines = [ln.strip() for ln in blob.split("\n") if ln.strip()]
        title_line = lines[0][:200] if lines else f"Summit slide {idx + 1}"
        body = lines[1:48] if len(lines) > 1 else [blob[:3500]]
        add_title_content_slide(prs, title_line, body)
        if len(prs.slides) >= MAX_TOTAL_SLIDES - 1:
            break

    add_thank_you_slide(prs)

    if len(prs.slides) > MAX_TOTAL_SLIDES:
        raise RuntimeError(
            f"Deck has {len(prs.slides)} slides; max is {MAX_TOTAL_SLIDES}. "
            "Shorten SUMMIT_APPEND_INDICES or template."
        )

    renumber_footers(prs)

    prs.save(str(OUT))
    print(f"Wrote {OUT}")
    print(f"Slides: {len(prs.slides)} (max {MAX_TOTAL_SLIDES})")


if __name__ == "__main__":
    merge()
