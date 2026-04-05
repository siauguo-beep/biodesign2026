#!/usr/bin/env python3
"""
Single canonical BDC deck + PDF for GitHub in-browser preview.
GitHub does NOT preview .pptx; use the generated .pdf on github.com/blob/.../file.pdf
"""

from pathlib import Path

from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
EXPORT = ROOT / "slides" / "export"
OUT_PPTX = EXPORT / "Archive_of_Extinction_Final_BDC2026.pptx"
OUT_PDF = EXPORT / "Archive_of_Extinction_Final_BDC2026.pdf"

# Final title (synthesized from your options + project research)
MAIN_TITLE = "Archive of Extinction"
SUBTITLE_LINES = [
    "Environmental Hypothesis Dossiers",
    "An AI Memorial for Lost Species - Sensory Time Capsule",
    "Biodesign Challenge 2026 | Convergent Life | Biodigital Excellence",
    "Macau University",
]

INSTITUTION = "Macau University"
MENTOR = "Atticus SIMS"
TEAM = [
    ("GUO XIAO YUE", "MC569254"),
    ("LIU JIA QUN", "MC569293"),
]


def ascii_pdf(text: str) -> str:
    """FPDF core fonts: Windows-1252-ish; keep ASCII for reliability."""
    return (
        text.replace("\u2014", "-")
        .replace("\u2013", "-")
        .replace("\u201c", '"')
        .replace("\u201d", '"')
        .replace("\u2019", "'")
        .replace("\u2192", "->")
        .replace("\u00b7", "-")
    )


def slide_plan():
    """Ordered content: (kind, title, body) where body is str or list[str] or dict for two_col."""
    s = []
    s.append(
        (
            "title",
            MAIN_TITLE,
            "\n".join(SUBTITLE_LINES),
        )
    )
    s.append(
        (
            "bullets",
            "Team & institution",
            [
                f"Institution: {INSTITUTION}",
                f"Mentor: {MENTOR}",
                f"Team: {TEAM[0][0]} - ID {TEAM[0][1]}",
                f"Team: {TEAM[1][0]} - ID {TEAM[1][1]}",
                "Program: MDes - Digital Art & AI Technology",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Title rationale (for judges)",
            [
                "Archive of Extinction: names the loss of species and of their temporal niches.",
                "Environmental Hypothesis Dossiers: each scene is a testable claim about Umwelt + habitat, tied to papers.",
                "AI memorial + sensory time capsule: generative WebXR makes rhythm and sensation discussable, not only images.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Biodesign Challenge 2026",
            [
                "Intersection of biotechnology, art, and design; Summit + MoMA judging for finalists.",
                'Gallery theme: "Convergent Life." We target Biodigital Excellence + Narrative, Context, Reflection.',
                "Judging: Narrative, Concept, Context, Reflection (each /4; strong tier ~2.75+).",
                "biodesignchallenge.org/judging | /prizes",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Core research proposition",
            [
                "If we use generative AI + WebXR to reconstruct and visualize temporal phenotypes of extinct species,",
                "can we repair the human gap in sensing extinction - the loss of another species' time, not only its image?",
                "Intersection statement from ideation workbook, stress-tested against the BDC rubric.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Conceptual methods",
            [
                "Palaeo-chronobiology: infer day/year phase from genomics, photoperiod proxies, morphology (orbit, channels).",
                "Temporal niche: extinction as erasure of shared ecological time, not only DNA.",
                "Umwelt (von Uexkull): model sensory worlds as hypotheses - never silent speculation.",
                "Red-flag test: if biology OR code is removed and the project still works, we failed biodesign or our digital edge.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Design techniques & stack",
            [
                "WebXR: A-Frame / Three.js; immersive dossier navigation.",
                "Generative environments: constrained by paleo proxies (Zimov steppe, seasonal light) - not fantasy biomes.",
                "Time scrubber: polar day/night, season; sonification of phase (optional).",
                "Epistemic UI: Cited / Interpolated / Speculative toggles + on-screen (Author, year) tags per scene ID.",
                "Ethics fork: resource sliders after Sherkow & Greely (2013); Indigenous context UI with Palawa-first citations.",
                "Reflection: templates/reflection-log-webxr.html - keyword match to ethics themes (heuristic, auditable).",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Ideation process (8 steps)",
            [
                "1 Map skills - 2 Intersection statements - 3 Rubric stress-test - 4 Biology literature sprints",
                "5 Bio/digital split - 6 Narrative draft - 7 Prototype + BDC deliverables - 8 Build + polish + Q&A",
                "Aligned with biodesign_cursor_agent.md + course workbook.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Design sprint (4 weeks)",
            [
                "W1 Discovery: citation grid, storyboard Scene_* IDs, WebXR shell.",
                "W2 Prototype: time axis, generative pipeline doc, epistemic wireframes.",
                "W3 Integration: sonification, UI_Ethics_Fork, UI_Indigenous_Context, reflection panel.",
                "W4 Polish: 10+5 min rehearsal, creative video, deploy, judge one-sheet.",
            ],
        )
    )
    s.append(
        (
            "two_col",
            "Two species - two dossiers",
            {
                "left_h": "Woolly mammoth",
                "left": [
                    "Arctic rhythm + mammoth steppe; scenes Scene_Bio_01 - Scene_Sensory_01.",
                    "Lynch 2015 circadian-related genes; Lu 2010 reindeer [Interpolated]; Zimov 2012 habitat.",
                    "TRPV3: Lynch 2015 primary; Kim 2017 bioRxiv optional + labeled.",
                ],
                "right_h": "Thylacine",
                "right": [
                    "Crepuscular / night vision metaphor; colonial context mandatory.",
                    "Pozniak 2018 orbit; Mass & Supin 2020 RGC [Interpolated]; Paddle; Sleightholme & Campbell.",
                    "T5: Clements + Palawa: Rimmer, Lehman; Schlunke 2025.",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            "Experience architecture (scene IDs)",
            [
                "Mammoth: Scene_Bio_01 -> Scene_Environment_02 -> Scene_Landscape_01 -> Scene_Sensory_01.",
                "Thylacine: Scene_Bio_02 -> Scene_POV_01 -> Scene_Context_01.",
                "Shared: UI_Indigenous_Context -> UI_Ethics_Fork -> UI_Reflection_Log.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Literature - Mammoth slots M1-M5",
            [
                "M1 Lynch et al. 2015 Cell Rep 10.1016/j.celrep.2015.06.027",
                "M2 Lu et al. 2010 Curr Biol 10.1016/j.cub.2010.01.042",
                "M3 Zimov et al. 2012 QSR 10.1016/j.quascirev.2012.08.004",
                "M4 Kim et al. 2017 bioRxiv 10.1101/151746 (preprint)",
                "M5 Sherkow & Greely 2013 Science 10.1126/science.1236946",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Literature - Thylacine T1-T5 + sovereignty",
            [
                "T1 Pozniak 2018 10.1002/ar.23910 | T2 Mass & Supin 2020 10.31857/S0002332920060107",
                "T3 Paddle 2000 ISBN 9780521782196 | T4 Sleightholme & Campbell 2016 10.1080/00222933.2016.1217037",
                "T5 Clements 2025 10.1177/13534858251272203",
                "Palawa: Rimmer 2020 Artlink | Lehman 2013 Griffith Review | Schlunke 2025 10.1017/ext.2025.10008",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "BDC rubric self-assessment",
            [
                "Narrative: discover rhythm -> loss -> ethical fork (target strong).",
                "Concept: time as spine; one biological proxy per scene.",
                "Context: de-extinction trade-offs + Indigenous sovereignty (expand citations).",
                "Reflection: interactive epistemic UI + user reflection + AI limits disclosed.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Deliverables & physical strategy",
            [
                "10 min talk + 5 min Q&A; 1-5 min creative video; slides to Drive.",
                "Visuals + WebXR; minimal tactile anchor + QR (no wet-lab for v1).",
                "Planning source of truth: BDC_2026_Extinction_Archive_Planning_Document.md + PROJECT_PLAN.md.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Risks we mitigate",
            [
                '"Just VR art" - always show citation strip.',
                "Overclaiming PER2 - UI uses circadian-related genes until PDF verified.",
                "Token TEK - Palawa-led sources + positionality statement.",
            ],
        )
    )
    s.append(
        (
            "title",
            "Thank you",
            f"{MAIN_TITLE}\nQuestions - biodesignchallenge.org",
        )
    )
    return s


def build_pptx(slides):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for item in slides:
        kind = item[0]
        if kind == "title":
            _, title, sub = item
            sl = prs.slides.add_slide(prs.slide_layouts[0])
            sl.shapes.title.text = title
            sl.placeholders[1].text = sub
        elif kind == "bullets":
            _, title, bullets = item
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = title
            tf = sl.placeholders[1].text_frame
            tf.clear()
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = Pt(16)
        elif kind == "two_col":
            _, title, d = item
            sl = prs.slides.add_slide(prs.slide_layouts[3])
            sl.shapes.title.text = title
            bodies = [p for p in sl.placeholders if p.placeholder_format.idx > 0]
            if len(bodies) < 2:
                flat = [d["left_h"]] + d["left"] + [""] + [d["right_h"]] + d["right"]
                sl = prs.slides.add_slide(prs.slide_layouts[1])
                sl.shapes.title.text = title
                tf = sl.placeholders[1].text_frame
                tf.clear()
                for i, line in enumerate(flat):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(14)
                continue
            tf_l = bodies[0].text_frame
            tf_l.clear()
            tf_l.text = d["left_h"]
            tf_l.paragraphs[0].font.bold = True
            tf_l.paragraphs[0].font.size = Pt(14)
            for b in d["left"]:
                p = tf_l.add_paragraph()
                p.text = b
                p.font.size = Pt(13)
            tf_r = bodies[1].text_frame
            tf_r.clear()
            tf_r.text = d["right_h"]
            tf_r.paragraphs[0].font.bold = True
            tf_r.paragraphs[0].font.size = Pt(14)
            for b in d["right"]:
                p = tf_r.add_paragraph()
                p.text = b
                p.font.size = Pt(13)

    EXPORT.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))


def build_pdf(slides):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_margins(left=18, top=18, right=18)

    def w():
        return pdf.epw

    for item in slides:
        kind = item[0]
        pdf.add_page()
        if kind == "title":
            _, title, sub = item
            pdf.set_font("Helvetica", "B", 22)
            pdf.multi_cell(w(), 10, ascii_pdf(title))
            pdf.ln(4)
            pdf.set_font("Helvetica", "", 12)
            pdf.multi_cell(w(), 6, ascii_pdf(sub))
        elif kind == "bullets":
            _, title, bullets = item
            pdf.set_font("Helvetica", "B", 16)
            pdf.multi_cell(w(), 8, ascii_pdf(title))
            pdf.ln(2)
            pdf.set_font("Helvetica", "", 11)
            for b in bullets:
                line = ascii_pdf("- " + b)
                pdf.multi_cell(w(), 5.5, line)
                pdf.ln(1)
        elif kind == "two_col":
            _, title, d = item
            pdf.set_font("Helvetica", "B", 16)
            pdf.multi_cell(w(), 8, ascii_pdf(title))
            pdf.ln(2)
            pdf.set_font("Helvetica", "B", 12)
            pdf.multi_cell(w(), 6, ascii_pdf(d["left_h"]))
            pdf.set_font("Helvetica", "", 10)
            for b in d["left"]:
                pdf.multi_cell(w(), 5, ascii_pdf("- " + b))
            pdf.ln(2)
            pdf.set_font("Helvetica", "B", 12)
            pdf.multi_cell(w(), 6, ascii_pdf(d["right_h"]))
            pdf.set_font("Helvetica", "", 10)
            for b in d["right"]:
                pdf.multi_cell(w(), 5, ascii_pdf("- " + b))

    pdf.output(str(OUT_PDF))


def main():
    slides = slide_plan()
    build_pptx(slides)
    build_pdf(slides)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {OUT_PDF} (open in Preview or view on GitHub.com)")


if __name__ == "__main__":
    main()
