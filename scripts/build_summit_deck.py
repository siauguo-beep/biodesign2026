#!/usr/bin/env python3
"""Generate BDC Summit deck for Extinction Archive (python-pptx) + GitHub Markdown mirror."""

from pathlib import Path
from urllib.parse import quote

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "slides" / "export" / "BDC_Summit_Extinction_Archive_2026.pptx"
OUT_ALIAS = ROOT / "slides" / "export" / "Extinction_Archive_Summit_2026.pptx"
MD_OUT = ROOT / "slides" / "SUMMIT_DECK_GITHUB_PREVIEW.md"

INSTITUTION = "Macau University"
MENTOR = "Atticus SIMS"
TEAM = [
    ("GUO XIAO YUE", "MC569254"),
    ("LIU JIA QUN", "MC569293"),
]

# Final public title: line 1 = cover title; subtitle carries memorial + capsule framing
MAIN_TITLE = "Extinction Archive: Umwelt Hypothesis Dossiers"
TITLE_SUB = (
    "AI memorial for lost species · sensory time capsule\n"
    "Biodesign Challenge 2026 · “Convergent Life”\n"
    f"{INSTITUTION}"
)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_bullets(prs, title, bullets, size_pt=17):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(size_pt)
    return slide


def add_two_col(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = title
    bodies = [p for p in slide.placeholders if p.placeholder_format.idx > 0]
    if len(bodies) < 2:
        return add_bullets(prs, title, left_bullets + [""] + right_bullets)
    tf_l = bodies[0].text_frame
    tf_l.clear()
    tf_l.text = left_title
    tf_l.paragraphs[0].font.bold = True
    tf_l.paragraphs[0].font.size = Pt(14)
    for b in left_bullets:
        p = tf_l.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(15)
    tf_r = bodies[1].text_frame
    tf_r.clear()
    tf_r.text = right_title
    tf_r.paragraphs[0].font.bold = True
    tf_r.paragraphs[0].font.size = Pt(14)
    for b in right_bullets:
        p = tf_r.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(15)
    return slide


def deck_spec():
    """Single source of truth: (kind, payload) for PPTX + Markdown."""
    s = []
    s.append(("title", {"title": MAIN_TITLE, "subtitle": TITLE_SUB}))
    s.append(
        (
            "bullets",
            {
                "title": "Team & institution",
                "bullets": [
                    f"Institution: {INSTITUTION}",
                    f"Mentor: {MENTOR}",
                    f"Team: {TEAM[0][0]} — ID {TEAM[0][1]}",
                    f"Team: {TEAM[1][0]} — ID {TEAM[1][1]}",
                    "Program: MDes — Digital Art & AI Technology.",
                ],
                "size": 20,
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Why this title?",
                "bullets": [
                    "“Umwelt Hypothesis Dossiers” = each species as a **folder of testable claims** (genes, orbits, climate), not a fantasy safari.",
                    "“AI memorial” = generative media **constrained** by literature — grief and wonder without faking science.",
                    "“Sensory time capsule” = we reconstruct **when** they lived (phase, season), not only **what** they looked like.",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Biodesign Challenge 2026",
                "bullets": [
                    "Premier student competition: biotechnology × art × design.",
                    "2026 gallery: “Convergent Life.” Targets: **Biodigital Excellence**; Narrative, Context, Reflection.",
                    "Judging: biodesignchallenge.org/judging",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Roles & build scope",
                "bullets": [
                    "Experience lead: narrative, ethics, palaeo-chronobiology translation.",
                    "Tech / AI lead: WebXR (A-Frame / Three.js), generative environments, epistemic UI.",
                    "Strategy: 重内容、轻装置 — Web-first; minimal tactile anchor + QR.",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Problem → provocation",
                "bullets": [
                    "Extinction erases **temporal niches** (daily / seasonal rhythms), not only DNA.",
                    "Public memory = flat images; we lack **felt time** of another species’ world.",
                    "Provocation: can biodigital design surface that loss and force honest **de-extinction** debate?",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Core design question",
                "bullets": [
                    '“If we use **generative AI + WebXR** to **reconstruct** a **temporal phenotype**, can we repair the human gap in **sensing extinction**?”',
                    "Palaeo-chronobiology + morphology (orbit, clock pathways, TRPV3…) = evidence ladder.",
                    "Umwelt = sensory world as **hypothesis**, disclosed layer-by-layer.",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Bio × digital quadrants (our mapping)",
                "bullets": [
                    "**SENSE / VISUALIZE:** photoperiod sliders, orbit-informed POV filters, habitat maps from Zimov et al.",
                    "**OPTIMIZE:** prompt + data rules so generative landscapes respect paleo proxies.",
                    "**SIMULATE:** phase / sonification of “a species day” — rhythmic, not decorative.",
                    "Red-flag: if biology OR code is removed and the piece still “works,” we failed the biodigital brief.",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Creative concept — two dossiers only",
                "bullets": [
                    "**Woolly mammoth:** Arctic circadian genomics (Lynch 2015), reindeer analogy (Lu 2010), mammoth steppe (Zimov 2012), TRPV3 cold (Lynch; Kim preprint optional).",
                    "**Thylacine:** skull / diel activity (Pozniak 2018), RGC metaphor (Mass & Supin 2020), history (Paddle; Sleightholme), sovereignty (Clements + Palawa-led Rimmer, Lehman; Schlunke).",
                    "Journey: discover rhythm → epistemic layers → ethics sliders → reflection log (HTML template in repo).",
                ],
            },
        )
    )
    s.append(
        (
            "two_col",
            {
                "title": "Scene IDs → citations",
                "left": (
                    "Mammoth path",
                    [
                        "Scene_Bio_01 — Lynch (circadian-related genes)",
                        "Scene_Environment_02 — Lu [Interpolated]",
                        "Scene_Landscape_01 — Zimov",
                        "Scene_Sensory_01 — TRPV3 / cold",
                    ],
                ),
                "right": (
                    "Thylacine path",
                    [
                        "Scene_Bio_02 — Pozniak / orbit",
                        "Scene_POV_01 — Mass & Supin [Interpolated]",
                        "Scene_Context_01 — Paddle + Sleightholme",
                        "UI_Indigenous_Context → UI_Ethics_Fork → UI_Reflection_Log",
                    ],
                ),
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Design techniques & stack",
                "bullets": [
                    "**WebXR:** A-Frame / Three.js; browser-based Summit demo.",
                    "**Generative AI:** 360° / environment art **bounded** by citation-backed parameters (no silent photoreal “truth”).",
                    "**Epistemic UI:** toggles for Cited / Interpolated / Speculative; on-screen (Author, year) tags.",
                    "**Sonification (optional):** tie audio to phase / season from literature.",
                    "**Reflection:** `templates/reflection-log-webxr.html` — keyword → ethics cards (Sherkow & Greely themes).",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Ideation method (8 steps)",
                "bullets": [
                    "1 Map interests & skills · 2 Intersection statements (“code × bio × problem”)",
                    "3 Stress-test vs BDC rubric (Narrative, Concept, Context, Reflection)",
                    "4 Biology sprint per species (30–45 min lit passes)",
                    "5 Bio / digital split + red-flag test",
                    "6 200-word narrative (Problem → Insight → Solution → Impact)",
                    "7 Prototype + deliverable owners · 8 Build sprint: Discovery → Prototype → Integration → Polish",
                    "Guided by `biodesign_cursor_agent.md` + course workbook.",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Four-week production spine",
                "bullets": [
                    "W1 Discovery — literature grid, storyboard ↔ DOI table, WebXR shell.",
                    "W2 Prototype — time scrubber, generative constraint doc, epistemic toggles.",
                    "W3 Integration — sonification, Indigenous + ethics UI, reflection panel.",
                    "W4 Polish — 10′ talk, 5′ Q&A drill, 1–5′ video, deploy, judge cheat sheet.",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Context — colonial violence & Country",
                "bullets": [
                    "Thylacine loss tied to colonial harm to Palawa peoples and lutruwita (Tasmania).",
                    "Sources: Clements (2025); Schlunke (2025) poetry / Palawa names; **Palawa-led** Rimmer (2020), Lehman (2013).",
                    "UI protocol: author positionality; no tourism copy as TEK.",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Ethics layer — Sherkow & Greely (2013)",
                "bullets": [
                    "Sliders: lab de-extinction vs **in situ** conservation funding.",
                    "Legal novelty of revived lineages → liability / ecosystem risk.",
                    "Grounded in Science DOI 10.1126/science.1236946 — not generic sci-fi branches.",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Key literature — Mammoth (DOIs)",
                "bullets": [
                    "M1 Lynch et al., 2015 — 10.1016/j.celrep.2015.06.027",
                    "M2 Lu et al., 2010 — 10.1016/j.cub.2010.01.042",
                    "M3 Zimov et al., 2012 — 10.1016/j.quascirev.2012.08.004",
                    "M4 Kim bioRxiv — 10.1101/151746 (TRPV3 structure; phenotype cite Lynch 2015)",
                    "M5 Sherkow & Greely — 10.1126/science.1236946",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Key literature — Thylacine & sovereignty",
                "bullets": [
                    "T1 Pozniak 2018 — 10.1002/ar.23910",
                    "T2 Mass & Supin 2020 — 10.31857/S0002332920060107",
                    "T3 Paddle 2000 — ISBN 9780521782196",
                    "T4 Sleightholme & Campbell 2016 — 10.1080/00222933.2016.1217037",
                    "T5 Clements 2025 — 10.1177/13534858251272203",
                    "Schlunke 2025 — 10.1017/ext.2025.10008 · Rimmer 2020 Artlink · Lehman 2013 Griffith Review",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Rubric self-check",
                "bullets": [
                    "Narrative: rhythm → loss → fork.",
                    "Concept: **time** as spine; one proxy per scene.",
                    "Context: trade-offs + sovereignty sources.",
                    "Reflection: epistemic UI + user reflection + AI limits.",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "BDC deliverables",
                "bullets": [
                    "10′ oral + 5′ Q&A · visuals / WebXR · minimal physical + QR",
                    "1–5′ creative video · slides on Drive · optional @biodesigned",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            {
                "title": "Repository pointers",
                "bullets": [
                    "Canonical plan: BDC_2026_Extinction_Archive_Planning_Document.md",
                    "Executive: PROJECT_PLAN.md · Reflection template: templates/reflection-log-webxr.html",
                    "GitHub preview this talk: slides/SUMMIT_DECK_GITHUB_PREVIEW.md (this file is auto-generated).",
                ],
            },
        )
    )
    s.append(
        (
            "title",
            {
                "title": "Thank you",
                "subtitle": f"{MAIN_TITLE}\n{INSTITUTION} · Biodigital Excellence · biodesignchallenge.org",
            },
        )
    )
    return s


def build_pptx(prs, spec):
    for kind, payload in spec:
        if kind == "title":
            add_title_slide(prs, payload["title"], payload["subtitle"])
        elif kind == "bullets":
            add_bullets(
                prs,
                payload["title"],
                payload["bullets"],
                payload.get("size", 17),
            )
        elif kind == "two_col":
            lt, lb = payload["left"]
            rt, rb = payload["right"]
            add_two_col(prs, payload["title"], lt, lb, rt, rb)


def build_markdown(spec, repo_placeholder="YOUR_USER/YOUR_REPO", branch="main"):
    lines = [
        "# Summit deck — GitHub preview mirror",
        "",
        f"**Project:** {MAIN_TITLE}  ",
        "**Subtitle framing:** AI memorial for lost species · sensory time capsule · BDC 2026 · Convergent Life  ",
        f"**Institution:** {INSTITUTION}  ",
        "",
        "> This Markdown file is **auto-generated** by `scripts/build_summit_deck.py` so you can **read the full talk on GitHub** without PowerPoint.",
        "",
        "## Browser preview of the real `.pptx` (after public push)",
        "",
        "GitHub does not render PPTX. After uploading to a **public** repo, use Microsoft’s viewer:",
        "",
        "1. Raw file URL pattern:",
        f"   `https://raw.githubusercontent.com/{repo_placeholder}/{branch}/slides/export/Extinction_Archive_Summit_2026.pptx`",
        "2. URL-encode that string, then open:",
        "   `https://view.officeapps.live.com/op/embed.aspx?src=` + encoded URL",
        "",
        "See **`docs/GITHUB_PPT_PREVIEW.md`** for step-by-step instructions.",
        "",
        "---",
        "",
    ]
    n = 0
    for kind, payload in spec:
        n += 1
        if kind == "title":
            lines.append(f"## Slide {n}: {payload['title']}")
            lines.append("")
            for para in payload["subtitle"].split("\n"):
                lines.append(para.strip())
                lines.append("")
        elif kind == "bullets":
            lines.append(f"## Slide {n}: {payload['title']}")
            lines.append("")
            for b in payload["bullets"]:
                lines.append(f"- {b}")
            lines.append("")
        elif kind == "two_col":
            lines.append(f"## Slide {n}: {payload['title']}")
            lines.append("")
            lt, lb = payload["left"]
            rt, rb = payload["right"]
            lines.append(f"### {lt}")
            for b in lb:
                lines.append(f"- {b}")
            lines.append("")
            lines.append(f"### {rt}")
            for b in rb:
                lines.append(f"- {b}")
            lines.append("")

    raw_example = f"https://raw.githubusercontent.com/{repo_placeholder}/{branch}/slides/export/Extinction_Archive_Summit_2026.pptx"
    enc = quote(raw_example, safe="")
    lines.append("---")
    lines.append("")
    lines.append("### One-line Office viewer (replace USER/REPO after you push)")
    lines.append("")
    lines.append(f"```\nhttps://view.officeapps.live.com/op/embed.aspx?src={enc}\n```")
    lines.append("")
    MD_OUT.parent.mkdir(parents=True, exist_ok=True)
    MD_OUT.write_text("\n".join(lines), encoding="utf-8")


def main():
    spec = deck_spec()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_pptx(prs, spec)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    prs.save(str(OUT_ALIAS))
    build_markdown(spec)
    print(f"Wrote {OUT}")
    print(f"Wrote {OUT_ALIAS} (byte-identical duplicate; shorter name)")
    print(f"Wrote {MD_OUT} (GitHub-friendly preview)")


if __name__ == "__main__":
    main()
